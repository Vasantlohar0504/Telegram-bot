import os
import tempfile
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import A4
from docx import Document
from PIL import Image

# ------------------ PPTX → PDF (Text + Images, Python-only) ------------------
def pptx_to_pdf(pptx_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    pdf_path = os.path.join(
        output_dir,
        os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
    )

    prs = Presentation(pptx_path)
    c = canvas.Canvas(pdf_path, pagesize=A4)
    pdf_width, pdf_height = A4
    margin = 50

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    for slide in prs.slides:
        y = pdf_height - margin

        # -------- TITLE --------
        if slide.shapes.title:
            c.setFont("Helvetica-Bold", 16)
            c.drawString(margin, y, slide.shapes.title.text)
            y -= 30

        # -------- SHAPES --------
        for shape in slide.shapes:

            # ---------- TEXT ----------
            if shape.has_text_frame:
                c.setFont("Helvetica", 12)
                for para in shape.text_frame.paragraphs:
                    if y < margin:
                        c.showPage()
                        y = pdf_height - margin
                    c.drawString(margin, y, para.text)
                    y -= 18

            # ---------- IMAGE ----------
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                    tmp.write(shape.image.blob)
                    img_path = tmp.name

                img = Image.open(img_path)
                img_width, img_height = img.size

                # Scale image to fit page width
                max_width = pdf_width - 2 * margin
                scale = min(max_width / img_width, 0.7)
                w = img_width * scale
                h = img_height * scale

                if y - h < margin:
                    c.showPage()
                    y = pdf_height - margin

                c.drawImage(
                    img_path,
                    margin,
                    y - h,
                    width=w,
                    height=h
                )

                y -= h + 20
                img.close()
                os.remove(img_path)

        c.showPage()

    c.save()
    return pdf_path


# ------------------ DOCX → PDF ------------------
def docx_to_pdf(docx_path, output_dir):
    os.makedirs(output_dir, exist_ok=True)
    pdf_path = os.path.join(
        output_dir,
        os.path.splitext(os.path.basename(docx_path))[0] + ".pdf"
    )

    doc = Document(docx_path)
    c = canvas.Canvas(pdf_path, pagesize=A4)
    width, height = A4
    y = height - 50

    for para in doc.paragraphs:
        if y < 50:
            c.showPage()
            y = height - 50
        c.drawString(50, y, para.text)
        y -= 15

    c.save()
    return pdf_path


# ------------------ Generic office_to_pdf ------------------
def office_to_pdf(input_file, output_dir):
    ext = os.path.splitext(input_file)[1].lower()
    if ext in [".pptx", ".ppt"]:
        return pptx_to_pdf(input_file, output_dir)
    elif ext == ".docx":
        return docx_to_pdf(input_file, output_dir)
    else:
        raise ValueError("Unsupported file type for office_to_pdf")
